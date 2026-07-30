#!/usr/bin/env python3
"""
Builds `ehlnofey-tier-ladders.pptx` - the talk-through deck for the tier ladders.

Audience assumption: the room knows nothing about Skyrim modding. Every mechanic is
built up from scratch, and the deck never uses a term it has not defined on an
earlier slide.

Target renderer is **LibreOffice Impress**, so this deliberately avoids anything
Impress renders differently to PowerPoint: no SmartArt, no gradients, no shadows,
no theme colours, no table styles. Everything is an explicit fill on an explicit
shape, and every table cell is coloured individually.

Run:  python build-deck.py          (needs python-pptx; see the repo CLAUDE.md)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# --------------------------------------------------------------------------- palette

BG     = RGBColor(0x0E, 0x11, 0x16)   # near-black, faintly blue
PANEL  = RGBColor(0x18, 0x1D, 0x27)   # card fill
PANEL2 = RGBColor(0x21, 0x28, 0x34)   # header rows, emphasis cards
LINE   = RGBColor(0x2E, 0x37, 0x45)
TEXT   = RGBColor(0xEC, 0xEA, 0xE5)
MUTED  = RGBColor(0x96, 0xA0, 0xAF)
GOLD   = RGBColor(0xD4, 0xA6, 0x3C)
ICE    = RGBColor(0x63, 0xB3, 0xD4)
RED    = RGBColor(0xCC, 0x5B, 0x44)
GREEN  = RGBColor(0x6F, 0xB0, 0x7F)

# T1..T7, cool -> hot
TIER = [
    RGBColor(0x5B, 0x8F, 0xA8),
    RGBColor(0x6F, 0xA9, 0x8E),
    RGBColor(0xA8, 0xA8, 0x5B),
    RGBColor(0xD4, 0xA6, 0x3C),
    RGBColor(0xD4, 0x82, 0x3C),
    RGBColor(0xC4, 0x55, 0x3C),
    RGBColor(0xA6, 0x3C, 0x63),
]

HEAD = "Segoe UI"
BODY = "Segoe UI"
MONO = "Consolas"

W, H = 13.333, 7.5
ML   = 0.80          # left margin
CW   = W - 2 * ML    # content width

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

_n = [0]


# ----------------------------------------------------------------- small helpers

def _noline(shape):
    shape.line.fill.background()
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    return shape


def rect(sl, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
    s = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return _noline(s)


def text(sl, x, y, w, h, runs, size=16, color=TEXT, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6, line=None):
    """`runs` is a string, or a list of (string, {overrides}) tuples, or a list of
    paragraphs where each paragraph is itself such a list."""
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if line:
            p.line_spacing = line
        bits = para if isinstance(para, list) else [(para, {})]
        for s, over in bits:
            r = p.add_run()
            r.text = s
            f = r.font
            f.name = over.get("font", font)
            f.size = Pt(over.get("size", size))
            f.bold = over.get("bold", bold)
            f.color.rgb = over.get("color", color)
    return box


def slide(title=None, kicker=None, chapter=False):
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, W, H, BG)
    _n[0] += 1
    if title:
        if kicker:
            text(sl, ML, 0.46, CW, 0.3, kicker.upper(), size=11.5, color=GOLD, bold=True)
        text(sl, ML, 0.76, CW, 0.66, title, size=30, bold=True, font=HEAD)
        rect(sl, ML, 1.44, 1.15, 0.045, GOLD)
    if not chapter:
        text(sl, W - ML - 1.0, H - 0.52, 1.0, 0.25, str(_n[0]),
             size=10.5, color=RGBColor(0x54, 0x5D, 0x6B), align=PP_ALIGN.RIGHT)
    return sl


def card(sl, x, y, w, h, head, body, accent=GOLD, fill=PANEL,
         hsize=15.5, bsize=12.5, pad=0.28):
    rect(sl, x, y, w, h, fill)
    rect(sl, x, y, 0.042, h, accent)
    text(sl, x + pad, y + pad - 0.03, w - 2 * pad, 0.34, head,
         size=hsize, bold=True, font=HEAD)
    if body:
        text(sl, x + pad, y + pad + 0.36, w - 2 * pad, h - pad - 0.5, body,
             size=bsize, color=MUTED, line=1.22)


def chip(sl, x, y, w, h, label, fill, color=BG, size=13, bold=True):
    rect(sl, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(sl, x, y + (h - 0.24) / 2 - 0.02, w, 0.26, label, size=size,
         color=color, bold=bold, align=PP_ALIGN.CENTER)


def arrow(sl, x, y, w=0.5, color=MUTED, size=22):
    text(sl, x, y, w, 0.4, "→", size=size, color=color,
         align=PP_ALIGN.CENTER, bold=True)


def table(sl, x, y, w, rows, col_w, row_h=0.42, head_h=0.44, size=12.5,
          head_fill=PANEL2, cell_fill=PANEL, alt_fill=None, colors=None,
          aligns=None, head_color=TEXT):
    """rows[0] is the header. `colors` optionally maps (row, col) -> RGBColor."""
    nr, nc = len(rows), len(rows[0])
    h = head_h + row_h * (nr - 1)
    shp = sl.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shp.table
    tbl.first_row = False
    tbl.horz_banding = False

    # drop the built-in table style so Impress does not paint its own borders
    tblPr = tbl._tbl.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}tblPr")
    if tblPr is not None:
        for child in list(tblPr):
            if child.tag.endswith("}tableStyleId"):
                tblPr.remove(child)

    for i, cwi in enumerate(col_w):
        tbl.columns[i].width = Inches(cwi)
    tbl.rows[0].height = Inches(head_h)
    for r in range(1, nr):
        tbl.rows[r].height = Inches(row_h)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.14)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.045)
            cell.margin_bottom = Inches(0.045)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = head_fill
            elif alt_fill is not None and r % 2 == 0:
                cell.fill.fore_color.rgb = alt_fill
            else:
                cell.fill.fore_color.rgb = cell_fill

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = (aligns[c] if aligns else PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = str(val)
            f = run.font
            f.name = BODY
            f.size = Pt(size if r else size - 0.5)
            f.bold = (r == 0)
            if r == 0:
                f.color.rgb = head_color
            elif colors and (r, c) in colors:
                f.color.rgb = colors[(r, c)]
            else:
                f.color.rgb = TEXT if c == 0 else MUTED
    return tbl


def note(sl, y, body, color=GOLD, h=0.62):
    """Full-width callout strip."""
    rect(sl, ML, y, CW, h, PANEL2)
    rect(sl, ML, y, 0.042, h, color)
    text(sl, ML + 0.3, y + 0.02, CW - 0.6, h - 0.04, body, size=13.5,
         color=TEXT, anchor=MSO_ANCHOR.MIDDLE, line=1.18)


# =========================================================================== 1. title

sl = slide(chapter=True)          # chapter=True suppresses the page number
rect(sl, 0, 0, 0.16, H, GOLD)
text(sl, 1.5, 2.25, 10.5, 0.4, "A SKYRIM SPECIAL EDITION PLUGIN",
     size=12.5, color=GOLD, bold=True)
text(sl, 1.5, 2.72, 10.5, 1.3, "Ehlnofey", size=68, bold=True, font=HEAD)
text(sl, 1.5, 4.05, 10.2, 0.9,
     "Deleveling Skyrim — and the tier ladders that have to replace the scaling",
     size=21, color=RGBColor(0xC3, 0xCA, 0xD4), line=1.25)
rect(sl, 1.5, 5.12, 2.2, 0.03, LINE)
text(sl, 1.5, 5.42, 10.5, 0.9,
     [[("No prior modding knowledge assumed. ", {"color": MUTED}),
       ("Every mechanic is built up from scratch.", {"color": MUTED})]],
     size=13.5)

# ================================================================ 2. the problem

sl = slide("Skyrim scales itself to you", "the problem")

BOXW, BOXH = 2.9, 0.92
def scene(y, lvl, enemy, tint):
    chip(sl, ML, y, 2.5, BOXH, "", PANEL, size=1)
    text(sl, ML, y + 0.16, 2.5, 0.3, "YOU", size=10.5, color=MUTED,
         align=PP_ALIGN.CENTER, bold=True)
    text(sl, ML, y + 0.44, 2.5, 0.35, f"Level {lvl}", size=17, bold=True,
         align=PP_ALIGN.CENTER)
    arrow(sl, ML + 2.65, y + 0.24)
    chip(sl, ML + 3.3, y, 3.5, BOXH, "", PANEL, size=1)
    text(sl, ML + 3.3, y + 0.16, 3.5, 0.3, "BLEAK FALLS BARROW", size=10.5,
         color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    text(sl, ML + 3.3, y + 0.44, 3.5, 0.35, "unchanged", size=17, bold=True,
         color=MUTED, align=PP_ALIGN.CENTER)
    arrow(sl, ML + 6.95, y + 0.24)
    chip(sl, ML + 7.6, y, 3.3, BOXH, "", PANEL, size=1)
    text(sl, ML + 7.6, y + 0.16, 3.3, 0.3, "WHAT SPAWNS", size=10.5, color=MUTED,
         align=PP_ALIGN.CENTER, bold=True)
    text(sl, ML + 7.6, y + 0.44, 3.3, 0.35, enemy, size=17, bold=True,
         color=tint, align=PP_ALIGN.CENTER)

scene(2.05, 5,  "Draugr, level 5",  ICE)
scene(3.25, 45, "Draugr, level 45", RED)

note(sl, 4.72,
     "Same dungeon, same corridor, same chest — and the same fight, forever. "
     "The world is a mirror held up to your character sheet.")
text(sl, ML, 5.62, CW, 1.1,
     "Vanilla does this almost everywhere: enemies, their gear, and the loot in the "
     "boss chest are all generated relative to your level at the moment you first "
     "walk in. Nothing in the world has a difficulty of its own.",
     size=14, color=MUTED, line=1.3)

# ==================================================================== 3. the cost

sl = slide("What that costs you", "the problem")
cw3, gap = (CW - 0.36) / 2, 0.36
card(sl, ML, 2.0, cw3, 1.5, "The map has no danger",
     "Every direction out of Riverwood is equally safe. There is nowhere you are "
     "not supposed to go yet, so there is nothing to earn the right to.", ICE)
card(sl, ML + cw3 + gap, 2.0, cw3, 1.5, "Exploring pays nothing",
     "The good sword is in the chest because you are level 40, not because the "
     "chest was hard to reach. Reward tracks the player, not the place.", GOLD)
card(sl, ML, 3.68, cw3, 1.5, "No curve, in either direction",
     "You never outgrow a threat and you are never overmatched by one. Difficulty "
     "is a flat line that follows you for the whole game.", RED)
card(sl, ML + cw3 + gap, 3.68, cw3, 1.5, "The fiction stops making sense",
     "A roadside bandit and an ancient Nordic Deathlord end up at the same level. "
     "The names still imply a hierarchy; the numbers no longer honour it.", GREEN)
note(sl, 5.5, "The mod's whole job: give difficulty and reward back to places and "
              "creatures, and take them away from the player's level.")

# =================================================================== 4. the rules

sl = slide("Three rules everything answers to", "the design")
for i, (n, t, b, c) in enumerate([
    ("01", "The world does not scale",
     "Every enemy and every hoard has a level set once, by hand. Nothing recalculates "
     "against the player, ever.", GOLD),
    ("02", "Danger is legible",
     "If walking somewhere can kill you, the world must have said so first. A fixed "
     "world is only fair if you can read it.", ICE),
    ("03", "Reward follows place",
     "Good loot exists because of where it is — a deep barrow, a dragon's hoard, "
     "a named boss. Never because you happened to be level 40.", GREEN),
]):
    y = 1.95 + i * 1.18
    rect(sl, ML, y, CW, 1.0, PANEL)
    rect(sl, ML, y, 0.042, 1.0, c)
    text(sl, ML + 0.3, y + 0.24, 0.7, 0.5, n, size=22, bold=True, color=c, font=HEAD)
    text(sl, ML + 1.1, y + 0.17, 3.6, 0.4, t, size=17, bold=True, font=HEAD)
    text(sl, ML + 4.8, y + 0.16, CW - 5.3, 0.7, b, size=12.5, color=MUTED, line=1.2)
note(sl, 5.6, "Rule 2 is the load-bearing one. Once the world stops scaling, "
              "legibility is the only protection a player has — and it is the "
              "reason the rest of this talk is about names as much as numbers.", ICE)

# ============================================================== 5. leveled lists

sl = slide("How Skyrim actually decides what spawns", "the mechanism")
text(sl, ML, 1.78, CW, 0.5,
     "A spawn point in a dungeon does not contain an enemy. It contains a "
     "menu — modders call it a leveled list.", size=16, color=TEXT, line=1.25)

rows = [["Gate", "Row on the menu"],
        ["1",  "Bandit"],
        ["5",  "Bandit Outlaw"],
        ["9",  "Bandit Thug"],
        ["14", "Bandit Highwayman"],
        ["19", "Bandit Plunderer"],
        ["25", "Bandit Marauder"]]
table(sl, ML, 2.45, 4.5, rows, [1.0, 3.5], row_h=0.4, size=12.5)

for i, (y, h, t, b) in enumerate([
    (2.42, 1.00, "Every row carries a level gate",
     "The number on the left is the player level at which that row unlocks."),
    (3.54, 1.38, "The engine hides what you have not earned",
     "It removes every row gated above your level. Then one switch on the list — "
     "\"calculate from all levels ≤ player\" — decides what happens next: pick at "
     "random from everything left, or take only the highest row still standing."),
    (5.04, 1.12, "So the gate is the scaling",
     "At level 3 one row is live, so you always meet a plain Bandit. At level 45 the "
     "list has opened up. Same spawn point, different world."),
]):
    card(sl, ML + 5.1, y, CW - 5.1, h, t, b,
         [GOLD, ICE, RED][i], hsize=14.5, bsize=12, pad=0.24)

note(sl, 6.30, "That one number is the entire scaling mechanism — and flattening it "
               "to 1 collapses both behaviours into one, because then the whole list "
               "is the highest row. Live, always, at every player level.", h=0.55)

sl.notes_slide.notes_text_frame.text = (
    "The switch is the flag 'CalculateFromAllLevelsLessThanOrEqualPlayer'.\n\n"
    "205 of the 264 leveled actor lists in the plugin carry it; 59 do not. So the "
    "'picks at random from everything under your level' description on this slide is "
    "the majority case, not the universal one — without the flag a level-45 player "
    "meets a Bandit Marauder every single time.\n\n"
    "Why it stops mattering: once every gate is 1, 'all rows at or below my level' and "
    "'only the highest qualifying row' select the same set — the entire list. So the "
    "flag was never audited and never needed to be.\n\n"
    "The example on screen, LCharBanditMelee1H (039CFC), does carry the flag."
)

# ================================================================== 6. the fix

sl = slide("The fix: flatten the gate, keep the pool", "the mechanism")

text(sl, ML, 1.82, 3.4, 0.32, "BEFORE — AS SHIPPED", size=11.5, color=MUTED, bold=True)
before = [["Gate", "Row"], ["1", "Bandit"], ["5", "Outlaw"], ["9", "Thug"],
          ["14", "Highwayman"], ["19", "Plunderer"], ["25", "Marauder"]]
table(sl, ML, 2.18, 3.4, before, [0.95, 2.45], row_h=0.38, size=12)

arrow(sl, ML + 3.52, 3.5, w=0.6, color=GOLD, size=30)

text(sl, ML + 4.25, 1.82, 3.4, 0.32, "AFTER — EHLNOFEY", size=11.5, color=GOLD, bold=True)
after = [["Gate", "Row"], ["1", "Runt"], ["1", "Outlaw"], ["1", "Outlaw"],
         ["1", "Outlaw"], ["1", "Outlaw"], ["1", "Thug"], ["1", "Thug"],
         ["1", "Thug"], ["1", "Highwayman"]]
table(sl, ML + 4.25, 2.18, 3.4, after, [0.95, 2.45], row_h=0.29, head_h=0.34, size=11)

card(sl, ML + 8.05, 2.18, CW - 8.05, 1.45, "Every gate becomes 1",
     "Nothing is ever hidden, so the player's level never enters the calculation.",
     GOLD, hsize=14, bsize=12, pad=0.24)
card(sl, ML + 8.05, 3.78, CW - 8.05, 1.45, "Repeats are the weighting",
     "The engine has no weight field. Listing Outlaw four times is how you make it "
     "four times as likely.", ICE, hsize=14, bsize=12, pad=0.24)

note(sl, 5.62, "Which rows you leave on the menu is now the entire design. "
               "Plunderer and Marauder are simply gone from this list — that "
               "decision, repeated across ~2,800 records, is the mod.")

# =============================================================== 7. tier ladder

sl = slide("The tier ladder", "the design")
text(sl, ML, 1.78, CW, 0.4,
     "Seven rungs. Every place, creature family and roster in the mod is assigned one.",
     size=15, color=MUTED)

rows = [["Tier", "Level", "Step", "What it means", "Where you meet it"],
        ["T1", "4",  "—",  "Nuisance", "Animal dens, the weakest camps"],
        ["T2", "8",  "+4",  "Local trouble", "Bandit camps, shipwrecks, minor caves"],
        ["T3", "14", "+6",  "Established threat", "Forts, Forsworn camps, ordinary crypts"],
        ["T4", "21", "+7",  "Serious", "Dwemer ruins, Falmer hives, hagraven nests"],
        ["T5", "30", "+9",  "Deep", "Dragon Priest lairs, deep Nordic ruins"],
        ["T6", "40", "+10", "Apex", "Solstheim's peak, Apocrypha, Volkihar's court"],
        ["T7", "50", "+10", "Capstone", "A short named list, nothing else"]]
colors = {}
for r in range(1, 8):
    colors[(r, 0)] = TIER[r - 1]
    colors[(r, 1)] = TIER[r - 1]
    colors[(r, 3)] = TEXT
table(sl, ML, 2.3, CW, rows, [0.95, 1.0, 0.85, 3.1, 5.83],
      row_h=0.415, size=13, colors=colors,
      aligns=[PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER,
              PP_ALIGN.LEFT, PP_ALIGN.LEFT])
note(sl, 5.92, "\"T4\" is shorthand for \"the level-21 rung\". A T4 creature is level 21 "
               "in the first dungeon you meet it and in the last one, forever.",
     h=0.55)

# ============================================================ 8. why the numbers

sl = slide("The numbers are not invented", "the design")
text(sl, ML, 1.78, CW, 0.4,
     "The ladder is lifted straight off the deepest ladder Bethesda already wrote — "
     "the draugr.", size=15, color=MUTED)

steps = [("Draugr", 1), ("Restless", 6), ("Wight", 13),
         ("Scourge", 21), ("Deathlord", 30), ("(ebony)", 40)]
bw = (CW - 5 * 0.22) / 6
for i, (nm, lv) in enumerate(steps):
    x = ML + i * (bw + 0.22)
    c = TIER[min(i + 1, 6)]
    rect(sl, x, 2.4, bw, 1.05, PANEL)
    rect(sl, x, 2.4, bw, 0.045, c)
    text(sl, x, 2.62, bw, 0.3, nm, size=13, bold=True, align=PP_ALIGN.CENTER)
    text(sl, x, 2.98, bw, 0.35, str(lv), size=20, bold=True, color=c,
         align=PP_ALIGN.CENTER)

cw8 = (CW - 0.34) / 2
card(sl, ML, 3.74, cw8, 1.45, "T4, T5 and T6 are exactly those gates",
     "21, 30 and 40 are not chosen numbers — they are where vanilla already puts "
     "Scourge, Deathlord and the ebony Deathlord.", GOLD, hsize=14.5, bsize=12)
card(sl, ML + cw8 + 0.34, 3.74, cw8, 1.45, "The steps widen on purpose",
     "+4, +6, +7, +9, +10. A rung is roughly a constant ratio of power, not a fixed "
     "number of levels.", ICE, hsize=14.5, bsize=12)
note(sl, 5.42, "The bottom three rungs are lifted (1→4, 6→8, 13→14) so a "
               "low-tier dungeon still has internal texture. At level 1 everything in a "
               "room resolves identically and the place feels uniform.")

# ============================================================== 9. the pivot

sl = slide("Tiers belong to creatures, not to rooms", "the pivot")
cw9 = (CW - 0.4) / 2
rect(sl, ML, 1.95, cw9, 2.5, PANEL)
rect(sl, ML, 1.95, 0.042, 2.5, RED)
text(sl, ML + 0.3, 2.18, cw9 - 0.6, 0.35, "The first plan — abandoned",
     size=16, bold=True, font=HEAD, color=RED)
text(sl, ML + 0.3, 2.62, cw9 - 0.6, 1.7,
     "Tier every dungeon instead, using Skyrim's own \"encounter zones\" — a "
     "per-location min and max level.\n\n"
     "It died on two measurements: encounter zones govern 0.3% of the outdoors, and "
     "no zone can reach the gear an enemy is wearing.",
     size=12.5, color=MUTED, line=1.25)

rect(sl, ML + cw9 + 0.4, 1.95, cw9, 2.5, PANEL)
rect(sl, ML + cw9 + 0.4, 1.95, 0.042, 2.5, GREEN)
text(sl, ML + cw9 + 0.7, 2.18, cw9 - 0.6, 0.35, "What shipped instead",
     size=16, bold=True, font=HEAD, color=GREEN)
text(sl, ML + cw9 + 0.7, 2.62, cw9 - 0.6, 1.7,
     "The tier moves onto the creature family. Flatten every menu, then choose which "
     "rows stay on it.\n\n"
     "A Draugr Scourge is level 21 in Bleak Falls Barrow, in Labyrinthian, and in a "
     "nameless barrow at the edge of the map.",
     size=12.5, color=MUTED, line=1.25)

note(sl, 4.72, "This trades something real away: a place can no longer be tiered "
               "independently of what is inside it. A cave is exactly as dangerous as "
               "the things Bethesda chose to put in it.", RED, h=0.75)
text(sl, ML, 5.75, CW, 0.9,
     "What it buys is that the tier now travels with the enemy — which is what "
     "makes the next slide possible.", size=14, color=MUTED)

# ============================================================ 10. names are tiers

sl = slide("Bethesda already wrote the difficulty labels", "legibility")
text(sl, ML, 1.78, CW, 0.4,
     "The draugr line is a lore hierarchy that happens to be perfectly ordered by power.",
     size=15, color=MUTED)

ladder = [("Draugr", 1, 0), ("Restless Draugr", 6, 1), ("Draugr Wight", 13, 2),
          ("Draugr Scourge", 21, 3), ("Draugr Deathlord", 30, 4)]
bw = (CW - 4 * 0.28) / 5
for i, (nm, lv, ti) in enumerate(ladder):
    x = ML + i * (bw + 0.28)
    c = TIER[ti + 1]
    rect(sl, x, 2.42, bw, 1.5, PANEL)
    rect(sl, x, 2.42, bw, 0.05, c)
    text(sl, x + 0.12, 2.68, bw - 0.24, 0.6, nm, size=14, bold=True,
         align=PP_ALIGN.CENTER, line=1.12)
    text(sl, x, 3.42, bw, 0.4, f"level {lv}", size=15, bold=True, color=c,
         align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(sl, x + bw + 0.01, 3.02, w=0.26, color=LINE, size=15)

note(sl, 4.2, "The player reads the tier off the enemy's name. No tutorial, no HUD, "
              "no wiki — the health bar already says how much trouble you are in.", ICE)
text(sl, ML, 5.08, CW, 1.3,
     "This is what makes rule 2 survivable. Under the abandoned plan the names were "
     "only correlated with power — you met a Wight because the dungeon happened "
     "to be tier 4. Once the tier lives on the creature, the name becomes the fact. "
     "The same trick works for Dremora (Churl → Caitiff → Kynval → "
     "Kynreeve → Markynaz → Valkynaz) and for Falmer, vampires and Forsworn.",
     size=13.5, color=MUTED, line=1.3)

# ============================================================== 11. bandit roster

sl = slide("A roster, worked end to end: bandits", "the design")
text(sl, ML, 1.76, CW, 0.4, "Vanilla ships six rungs. The roster keeps four of them, "
     "weighted, and cuts two.", size=14.5, color=MUTED)

picks = [("Runt", 1, 1, True), ("Outlaw", 5, 4, True), ("Thug", 9, 3, True),
         ("Highwayman", 14, 1, True), ("Plunderer", 19, 0, False),
         ("Marauder", 25, 0, False)]
bw = (CW - 5 * 0.2) / 6
for i, (nm, lv, wt, keep) in enumerate(picks):
    x = ML + i * (bw + 0.2)
    c = TIER[min(i, 6)] if keep else RGBColor(0x3A, 0x42, 0x50)
    rect(sl, x, 2.3, bw, 1.42, PANEL if keep else RGBColor(0x13, 0x17, 0x1E))
    rect(sl, x, 2.3, bw, 0.045, c)
    text(sl, x, 2.52, bw, 0.3, nm, size=13, bold=True,
         color=TEXT if keep else RGBColor(0x5A, 0x63, 0x71), align=PP_ALIGN.CENTER)
    text(sl, x, 2.85, bw, 0.28, f"level {lv}", size=11.5,
         color=MUTED if keep else RGBColor(0x4A, 0x53, 0x61), align=PP_ALIGN.CENTER)
    text(sl, x, 3.2, bw, 0.4, (f"×{wt}" if keep else "cut"),
         size=17 if keep else 13, bold=True,
         color=c if keep else RGBColor(0x5A, 0x63, 0x71), align=PP_ALIGN.CENTER)

cwb = (CW - 0.34) / 2
card(sl, ML, 4.0, cwb, 1.24, "The camp sits at T2",
     "Mass on Outlaw and Thug. The weakest and strongest rungs are equally rare — "
     "1 in 9 each.", TIER[1], hsize=14.5, bsize=12)
card(sl, ML + cwb + 0.34, 4.0, cwb, 1.24, "The chief is pinned to level 28",
     "One rung, no spread — the top of the vanilla ladder. A camp's boss is the "
     "only fight in it that can carry real difficulty.", TIER[4], hsize=14.5, bsize=12)
note(sl, 5.48, "Nothing here writes a level onto a creature. Every roster is a filter "
               "and a weighting of rows Bethesda already wrote — which is why it "
               "stays compatible with almost everything.")

# =============================================================== 12. the trap

sl = slide("The trap: a spread of levels needs a spread of names", "found in play")
text(sl, ML, 1.76, CW, 0.62,
     "The chief was pinned to one level for a reason. Playtesting found it rolling "
     "anywhere from level 6 to level 28 — and every single one of those was "
     "called \"Bandit Chief\".", size=14.5, color=MUTED, line=1.25)

card(sl, ML, 2.52, CW, 1.24, "Why the name does not follow the level",
     "A creature record inherits its name and its level through two different "
     "switches. Set the level switch but not the name switch, and the name falls "
     "through to one shared record — so every power level shows the same string.",
     RED, hsize=15, bsize=12.5)

rows = [["Boss family", "Names across its rungs", "Verdict"],
        ["Draugr",   "Overlord · Wight Lord · Scourge Lord · Death Overlord", "spread is fine"],
        ["Falmer",   "Skulker · Gloomlurker · Nightprowler · Shadowmaster", "spread is fine"],
        ["Bandit",   "one name for every rung", "pinned to 28"],
        ["Forsworn", "one name for every rung", "open"],
        ["Warlock",  "one name for every rung", "open"],
        ["Thalmor",  "\"Thalmor Wizard\" at all seven rungs", "open"],
        ["Vampire",  "one name for every rung", "open"]]
vc = {}
for r, col in [(1, GREEN), (2, GREEN), (3, GOLD), (4, RED), (5, RED), (6, RED), (7, RED)]:
    vc[(r, 2)] = col
table(sl, ML, 3.9, CW, rows, [2.3, 6.9, 2.53], row_h=0.30, head_h=0.38, size=12,
      colors=vc)
note(sl, 6.5, "The rule that came out of it: a band of levels is only allowed where "
              "the rungs have different names.", ICE, h=0.5)

# ============================================================== 13. the wildlife

sl = slide("Wildlife: the map draws itself", "legibility")
text(sl, ML, 1.78, CW, 0.4,
     "Vanilla already partitions animal spawns by biome. Freeze each biome at its own "
     "tier and a gradient appears for free.", size=14.5, color=MUTED)

biomes = [("Plains", "T2", 1, "Wolf ×4  SabreCat ×1"),
          ("Forest · Canyon\nMarsh · Hills · Coast", "T3", 2,
           "Bear ×5  Spider ×4\nWolf ×4  Troll ×1"),
          ("Snowfields", "T4", 3, "WolfIce ×3  SabreCatSnow ×2\nBearSnow ×1"),
          ("Mountains", "T5", 4, "IceWraith ×4  WolfIce ×4\nFrost Troll ×3")]
bw = (CW - 3 * 0.3) / 4
for i, (nm, tier, ti, mix) in enumerate(biomes):
    x = ML + i * (bw + 0.3)
    c = TIER[ti]
    rect(sl, x, 2.4, bw, 2.15, PANEL)
    rect(sl, x, 2.4, bw, 0.05, c)
    text(sl, x + 0.16, 2.62, bw - 0.32, 0.62, nm, size=14, bold=True,
         align=PP_ALIGN.CENTER, line=1.1, anchor=MSO_ANCHOR.MIDDLE)
    text(sl, x, 3.3, bw, 0.34, tier, size=19, bold=True, color=c, align=PP_ALIGN.CENTER)
    text(sl, x + 0.14, 3.72, bw - 0.28, 0.7, mix, size=11, color=MUTED,
         align=PP_ALIGN.CENTER, line=1.25, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        arrow(sl, x + bw + 0.02, 3.28, w=0.28, color=LINE, size=16)

note(sl, 4.86, "Only the T5 mountain roster contains the Frost Troll. The mountain is "
               "the warning — and it is visible from anywhere in Skyrim.", ICE)
text(sl, ML, 5.72, CW, 0.9,
     "This matters because the mod switches off vanilla's twelve hidden \"do not spawn "
     "this yet\" flags. Geography has to carry the warning instead, and here it does.",
     size=13.5, color=MUTED, line=1.28)

# ============================================================= 14. what it costs

sl = slide("What it costs — honestly", "trade-offs")
cwc = (CW - 0.36) / 2
card(sl, ML, 1.94, cwc, 1.5, "Bandit camps go bimodal",
     "Trivial mooks at T2, then a level-28 chief. That is a spike, and it is a "
     "deliberate one — a camp has nothing else in it to carry difficulty.", GOLD)
card(sl, ML + cwc + 0.36, 1.94, cwc, 1.5, "Some content is simply shut",
     "A level-5 character cannot clear a deep barrow. That is the point, but it puts "
     "all the weight on signposting getting it right.", RED)
card(sl, ML, 3.6, cwc, 1.5, "Followers have to be hand-set",
     "They do not scale either, so each one needs a level chosen from role and lore. "
     "Picking a companion becomes a real decision.", ICE)
card(sl, ML + cwc + 0.36, 3.6, cwc, 1.5, "Level is not the same as threat",
     "A fixed world can be a dull one. Perks, spells and combat style do as much work "
     "as the number — and they are next.", GREEN)
note(sl, 5.36, "None of these are bugs to be fixed later. They are the shape a world "
              "takes once it stops flattering the player, and the design says so out "
              "loud rather than discovering them in testing.")

# ============================================================== 15. where it is

sl = slide("Where it stands", "status")
stats = [("2,833", "records, all overrides"),
         ("0", "new form IDs"),
         ("16,308", "list entries, every gate flat"),
         ("4", "master files")]
bw = (CW - 3 * 0.3) / 4
for i, (big, small) in enumerate(stats):
    x = ML + i * (bw + 0.3)
    rect(sl, x, 1.98, bw, 1.25, PANEL)
    rect(sl, x, 1.98, bw, 0.045, GOLD)
    text(sl, x, 2.28, bw, 0.5, big, size=30, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, font=HEAD)
    text(sl, x + 0.12, 2.84, bw - 0.24, 0.4, small, size=11.5, color=MUTED,
         align=PP_ALIGN.CENTER, line=1.15)

cwd = (CW - 0.36) / 2
card(sl, ML, 3.56, cwd, 1.62, "Working",
     "Builds and round-trips cleanly. In play at character level 1 and level 45, Bleak "
     "Falls Barrow and Swindler's Den give the same spread of enemies both times — "
     "the first real evidence rule 1 holds.", GREEN)
card(sl, ML + cwd + 0.36, 3.56, cwd, 1.62, "Still open",
     "Four boss families still fail the naming test. Followers are unset. Boss-chest "
     "loot has not been verified across two player levels yet.", RED)
note(sl, 5.42, "Everything shipped is an override of a record Bethesda wrote. No new "
               "records, so it stays as compatible as a mod of this reach can be.")

# ================================================================= 16. takeaways

sl = slide("Three things to take away", "close")
for i, (n, t, b) in enumerate([
    ("01", "The gate is the scaling",
     "One number on each row of a leveled list is what ties the whole world to your "
     "character sheet. Flatten it and the coupling is gone."),
    ("02", "Then the curation is the design",
     "Once nothing is hidden, which rows remain on the menu is the only lever left — "
     "and that is a design job, not a technical one."),
    ("03", "Names carry the difficulty",
     "A fixed world is only fair if it can be read, and Bethesda already wrote the "
     "labels. Where the labels run out, the spread has to stop."),
]):
    y = 1.98 + i * 1.32
    rect(sl, ML, y, CW, 1.14, PANEL)
    rect(sl, ML, y, 0.042, 1.14, [GOLD, ICE, GREEN][i])
    text(sl, ML + 0.32, y + 0.3, 0.8, 0.5, n, size=24, bold=True,
         color=[GOLD, ICE, GREEN][i], font=HEAD)
    text(sl, ML + 1.25, y + 0.2, 4.2, 0.4, t, size=17, bold=True, font=HEAD)
    text(sl, ML + 5.6, y + 0.2, CW - 6.1, 0.85, b, size=12.5, color=MUTED, line=1.22)

text(sl, ML, 6.1, CW, 0.5, "Questions", size=24, bold=True, color=GOLD, font=HEAD)

# ---------------------------------------------------------------------- write

prs.core_properties.title = "Ehlnofey - Deleveling Skyrim"
prs.core_properties.subject = "Tier ladders"
prs.core_properties.author = "Ehlnofey"

out = Path(__file__).with_name("ehlnofey-tier-ladders.pptx")
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
